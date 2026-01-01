"""PPTX export functionality.

This module provides:
- PPTXExporter: Convert HTML slides to PowerPoint format

Uses Selenium with system Chrome/Chromium for screenshots, which automatically
inherits system fonts (including Chinese fonts).

Supports parallel screenshot capture for faster export.
"""

import asyncio
import logging
import re
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# Check for Selenium availability
try:
    from selenium import webdriver
    from selenium.webdriver.chrome.service import Service as ChromeService
    from selenium.webdriver.chrome.options import Options as ChromeOptions
    from webdriver_manager.chrome import ChromeDriverManager
    SELENIUM_AVAILABLE = True
except ImportError:
    SELENIUM_AVAILABLE = False


class PPTXExporter:
    """Exports slides as a PowerPoint presentation.
    
    Features:
    - Converts HTML slides to PPTX
    - Takes screenshots of HTML using Selenium (inherits system fonts)
    - Adds speaker notes
    """
    
    def __init__(self, workspace_path: Path):
        """Initialize exporter.
        
        Args:
            workspace_path: Path to task workspace
        """
        self.workspace_path = Path(workspace_path)
        self.slides_path = self.workspace_path / "slides"
        self.screenshots_path = self.workspace_path / "screenshots"
    
    def get_slide_files(self) -> List[Path]:
        """Get list of slide HTML files in order."""
        if not self.slides_path.exists():
            return []
        
        files = list(self.slides_path.glob("slide_*.html"))
        
        def get_slide_num(path: Path) -> int:
            match = re.search(r'slide_(\d+)\.html', path.name)
            return int(match.group(1)) if match else 0
        
        return sorted(files, key=get_slide_num)
    
    def _create_chrome_options(self, width: int, height: int) -> 'ChromeOptions':
        """Create Chrome options for headless screenshot."""
        chrome_options = ChromeOptions()
        chrome_options.add_argument("--headless=new")
        chrome_options.add_argument("--disable-gpu")
        chrome_options.add_argument("--no-sandbox")
        chrome_options.add_argument("--disable-dev-shm-usage")
        chrome_options.add_argument(f"--window-size={width},{height}")
        # Force device scale factor for higher quality
        chrome_options.add_argument("--force-device-scale-factor=2")
        return chrome_options
    
    def _take_single_screenshot(
        self,
        slide_path: Path,
        screenshot_path: Path,
        width: int,
        height: int,
        driver_path: str
    ) -> Tuple[Path, bool, str]:
        """Take a screenshot of a single slide.
        
        Args:
            slide_path: Path to the HTML slide
            screenshot_path: Path to save the screenshot
            width: Viewport width
            height: Viewport height
            driver_path: Path to ChromeDriver
        
        Returns:
            Tuple of (screenshot_path, success, error_message)
        """
        driver = None
        try:
            chrome_options = self._create_chrome_options(width, height)
            service = ChromeService(driver_path)
            driver = webdriver.Chrome(service=service, options=chrome_options)
            driver.set_window_size(width, height)
            
            # Navigate to the slide file
            file_url = f"file:///{slide_path.resolve().as_posix()}"
            driver.get(file_url)
            
            # Wait for page to load and web fonts
            time.sleep(1.5)
            
            # Try to wait for fonts to be ready
            try:
                driver.execute_script("return document.fonts.ready")
            except Exception:
                pass  # Ignore if fonts API not available
            
            # Take screenshot
            driver.save_screenshot(str(screenshot_path))
            logger.info(f"Screenshot saved: {screenshot_path}")
            return (screenshot_path, True, "")
            
        except Exception as e:
            logger.error(f"Failed to screenshot {slide_path}: {e}")
            return (screenshot_path, False, str(e))
            
        finally:
            if driver:
                driver.quit()
    
    def take_screenshots_sync(
        self,
        width: int = 1920,
        height: int = 1080,
        max_workers: int = 4
    ) -> List[Path]:
        """Take screenshots of all slides using Selenium in parallel.
        
        Uses system Chrome/Chromium which inherits system fonts,
        avoiding Chinese font display issues.
        
        Args:
            width: Viewport width (default: 1920)
            height: Viewport height (default: 1080)
            max_workers: Maximum parallel Chrome instances (default: 4)
        
        Returns:
            List of screenshot paths
        
        Raises:
            RuntimeError: If Selenium is not available
        """
        if not SELENIUM_AVAILABLE:
            raise RuntimeError(
                "Selenium is not installed. Please run:\n"
                "  pip install selenium webdriver-manager"
            )
        
        slides = self.get_slide_files()
        if not slides:
            return []
        
        self.screenshots_path.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"Taking screenshots of {len(slides)} slides in parallel (max {max_workers} workers)...")
        
        # Pre-install ChromeDriver once (shared by all workers)
        driver_path = ChromeDriverManager().install()
        
        # Prepare tasks
        tasks = [
            (slide_path, self.screenshots_path / f"{slide_path.stem}.png")
            for slide_path in slides
        ]
        
        # Execute in parallel
        screenshot_results = []
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(
                    self._take_single_screenshot,
                    slide_path,
                    screenshot_path,
                    width,
                    height,
                    driver_path
                ): (slide_path, screenshot_path)
                for slide_path, screenshot_path in tasks
            }
            
            for future in as_completed(futures):
                screenshot_path, success, error = future.result()
                if success:
                    screenshot_results.append(screenshot_path)
                else:
                    logger.warning(f"Screenshot failed: {screenshot_path}")
        
        # Sort results by slide number
        def get_slide_num(path: Path) -> int:
            match = re.search(r'slide_(\d+)', path.stem)
            return int(match.group(1)) if match else 0
        
        screenshot_results.sort(key=get_slide_num)
        
        logger.info(f"Created {len(screenshot_results)} screenshots")
        return screenshot_results
    
    async def take_screenshots(self) -> List[Path]:
        """Take screenshots of all slides (async wrapper).
        
        Returns:
            List of screenshot paths
        """
        # Run sync version in thread pool
        return await asyncio.to_thread(self.take_screenshots_sync)
    
    async def export(
        self,
        output_path: Optional[Path] = None,
        title: str = "Presentation",
        take_screenshots: bool = True
    ) -> Path:
        """Export slides as PPTX.
        
        Args:
            output_path: Output file path
            title: Presentation title
            take_screenshots: Whether to take screenshots for slides
            
        Returns:
            Path to exported PPTX file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError(
                "python-pptx not installed. Run: pip install python-pptx"
            )
        
        if output_path is None:
            output_path = self.workspace_path / "presentation.pptx"
        
        slides = self.get_slide_files()
        if not slides:
            raise ValueError("No slides found to export")
        
        # Take screenshots if needed
        if take_screenshots:
            screenshot_paths = await self.take_screenshots()
        else:
            screenshot_paths = list(self.screenshots_path.glob("slide_*.png"))
            screenshot_paths.sort(key=lambda p: int(re.search(r'slide_(\d+)', p.stem).group(1)))
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 widescreen
        prs.slide_height = Inches(7.5)
        
        # Get blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        
        for i, slide_path in enumerate(slides):
            slide = prs.slides.add_slide(blank_layout)
            
            # Add screenshot as background if available
            if i < len(screenshot_paths) and screenshot_paths[i].exists():
                slide.shapes.add_picture(
                    str(screenshot_paths[i]),
                    Inches(0), Inches(0),
                    width=Inches(13.333), height=Inches(7.5)
                )
            
            # Add speaker notes if available
            notes = self._extract_notes(slide_path)
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes
        
        # Save presentation
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        return output_path
    
    def _extract_notes(self, slide_path: Path) -> str:
        """Extract speaker notes from slide HTML.
        
        Looks for notes in data attributes or comments.
        """
        content = slide_path.read_text(encoding="utf-8")
        
        # Try to find notes in data attribute
        match = re.search(r'data-notes="([^"]*)"', content)
        if match:
            return match.group(1)
        
        # Try to find notes in HTML comment
        match = re.search(r'<!--\s*notes:\s*(.*?)\s*-->', content, re.DOTALL | re.IGNORECASE)
        if match:
            return match.group(1).strip()
        
        return ""
    
    def export_sync(
        self,
        output_path: Optional[Path] = None,
        title: str = "Presentation"
    ) -> Path:
        """Synchronous wrapper for export."""
        return asyncio.run(self.export(output_path, title))