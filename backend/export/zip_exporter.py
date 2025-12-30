"""ZIP package export functionality.

This module provides:
- ZipExporter: Package all presentation assets into a ZIP file
"""

import io
import json
import logging
import zipfile
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Union

from config import Settings, get_settings

logger = logging.getLogger(__name__)


class ZipExporter:
    """Packages presentation assets into a ZIP file.
    
    Includes:
    - HTML slides
    - Combined presentation (iframe loader)
    - PPTX file (via screenshots)
    - Speech scripts and coaching
    - Presentation plan
    - README with keyboard shortcuts
    """
    
    def __init__(
        self,
        workspace_path: Path,
        settings: Optional[Settings] = None
    ):
        """Initialize exporter.
        
        Args:
            workspace_path: Path to task workspace
            settings: Application settings
        """
        self.workspace_path = Path(workspace_path)
        self.slides_path = self.workspace_path / "slides"
        self.settings = settings or get_settings()
    
    def _get_manifest(self) -> Optional[dict]:
        """Read and return the manifest.json file."""
        manifest_path = self.slides_path / "manifest.json"
        if not manifest_path.exists():
            return None
        try:
            return json.loads(manifest_path.read_text(encoding='utf-8'))
        except Exception:
            return None
    
    def _screenshot_slides(self) -> List[Path]:
        """Take screenshots of all slide HTML files using Selenium.
        
        Returns:
            List of paths to screenshot files
        """
        try:
            from selenium import webdriver
            from selenium.webdriver.chrome.service import Service as ChromeService
            from selenium.webdriver.chrome.options import Options as ChromeOptions
            from webdriver_manager.chrome import ChromeDriverManager
        except ImportError:
            logger.warning(
                "Selenium not installed. Run: pip install selenium webdriver-manager"
            )
            return []
        
        # Determine output directory
        output_dir = self.workspace_path / "screenshots"
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Read manifest to get slide order
        manifest = self._get_manifest()
        if not manifest:
            logger.warning("Manifest not found, cannot take screenshots")
            return []
        
        slides_meta = manifest.get("slides", [])
        screenshot_paths = []
        
        width = 1920
        height = 1080
        
        logger.info(f"Taking screenshots of {len(slides_meta)} slides using Selenium...")
        
        # Configure Chrome options for headless mode
        chrome_options = ChromeOptions()
        chrome_options.add_argument("--headless=new")
        chrome_options.add_argument("--disable-gpu")
        chrome_options.add_argument("--no-sandbox")
        chrome_options.add_argument("--disable-dev-shm-usage")
        chrome_options.add_argument(f"--window-size={width},{height}")
        # Force device scale factor for higher quality
        chrome_options.add_argument("--force-device-scale-factor=2")
        
        driver = None
        try:
            # Auto-install ChromeDriver
            service = ChromeService(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=chrome_options)
            driver.set_window_size(width, height)
            
            import time
            
            for i, slide_meta in enumerate(slides_meta):
                slide_file = slide_meta.get("file", f"slide_{i+1}.html")
                slide_path = self.slides_path / slide_file
                
                if not slide_path.exists():
                    logger.warning(f"Slide file not found: {slide_path}")
                    continue
                
                # Navigate to the slide file
                file_url = f"file:///{slide_path.resolve().as_posix()}"
                driver.get(file_url)
                
                # Wait a bit for any animations/transitions to complete
                time.sleep(0.8)
                
                # Take screenshot
                screenshot_name = f"slide_{i+1:03d}.png"
                screenshot_path = output_dir / screenshot_name
                driver.save_screenshot(str(screenshot_path))
                
                screenshot_paths.append(screenshot_path)
                logger.info(f"Screenshot saved: {screenshot_path}")
                
        except Exception as e:
            logger.error(f"Error taking screenshots: {e}")
        finally:
            if driver:
                driver.quit()
        
        logger.info(f"Created {len(screenshot_paths)} screenshots in {output_dir}")
        return screenshot_paths
    
    def _create_pptx_from_screenshots(
        self,
        screenshot_paths: List[Path],
        output_path: Path,
        title: str = "Presentation"
    ) -> Optional[Path]:
        """Create a PowerPoint file from slide screenshots.
        
        Args:
            screenshot_paths: List of paths to screenshot images
            output_path: Output path for the PPTX file
            title: Title for the presentation
        
        Returns:
            Path to the created PPTX file, or None if failed
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            logger.warning(
                "python-pptx not installed. Run: pip install python-pptx"
            )
            return None
        
        try:
            # Create presentation with 16:9 aspect ratio
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9 at standard size
            prs.slide_height = Inches(7.5)
            
            # Blank slide layout (index 6 is typically blank)
            blank_layout = prs.slide_layouts[6]
            
            for screenshot_path in screenshot_paths:
                if not screenshot_path.exists():
                    logger.warning(f"Screenshot not found: {screenshot_path}")
                    continue
                
                # Add a blank slide
                slide = prs.slides.add_slide(blank_layout)
                
                # Add the screenshot as a full-slide background image
                slide.shapes.add_picture(
                    str(screenshot_path),
                    Inches(0),
                    Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )
            
            # Save the presentation
            prs.save(str(output_path))
            logger.info(f"Created PPTX: {output_path}")
            
            return output_path
        except Exception as e:
            logger.error(f"Error creating PPTX: {e}")
            return None
    
    def _generate_speech_materials(
        self
    ) -> tuple[Optional[str], Optional[str]]:
        """Generate speech script and coaching.
        
        Returns:
            Tuple of (speech_script, speech_coaching) or (None, None) if failed
        """
        try:
            from export.speech_generator import SpeechGenerator
        except ImportError:
            logger.warning("SpeechGenerator not available")
            return None, None
        
        # Check if API key is configured (SpeechGenerator reads from .env via settings)
        if not self.settings.openai_api_key:
            logger.warning("OPENAI_API_KEY not configured in .env, skipping speech generation")
            return None, None
        
        try:
            # SpeechGenerator reads configuration from .env via get_settings()
            generator = SpeechGenerator(settings=self.settings)
            
            speech_script = generator.generate_speech_script(self.slides_path)
            speech_coaching = None
            
            if speech_script:
                speech_coaching = generator.generate_speech_coaching(speech_script)
            
            return speech_script, speech_coaching
        except Exception as e:
            logger.error(f"Error generating speech materials: {e}")
            return None, None
    
    def _generate_readme(
        self,
        title: str,
        has_pptx: bool = False,
        has_speech: bool = False,
        has_coaching: bool = False
    ) -> str:
        """Generate README content for the ZIP package."""
        pptx_note = ""
        if has_pptx:
            pptx_note = "├── presentation.pptx    # PowerPoint 版本\n"
        
        speech_note = ""
        if has_speech:
            speech_note = "├── speech_script.md     # 演讲稿\n"
        if has_coaching:
            speech_note += "├── speech_coaching.md   # 演讲指导\n"
        
        return f"""# {title}

## 使用说明

打开 `presentation.html` 文件即可开始演示。
{"也可以使用 `presentation.pptx` 在 PowerPoint 中打开。" if has_pptx else ""}

{"## 演讲辅助材料" if has_speech else ""}
{"- `speech_script.md` - 完整的演讲稿，包含每页幻灯片的演讲内容和时间估算" if has_speech else ""}
{"- `speech_coaching.md` - 演讲技巧指导，帮助您更好地完成演讲" if has_coaching else ""}

## 键盘快捷键 (HTML 版本)

| 按键 | 功能 |
|------|------|
| ← | 上一页 |
| → | 下一页 |
| 空格 | 下一页 |
| Enter | 切换全屏 |
| F | 切换全屏 |
| Esc | 退出全屏 |
| Home | 跳到第一页 |
| End | 跳到最后一页 |

## 文件结构

```
├── presentation.html    # 主文件，打开此文件开始演示
{pptx_note}{speech_note}├── slides/              # 幻灯片文件目录
│   ├── slide_1.html
│   ├── slide_2.html
│   └── ...
└── README.md            # 本说明文件
```

## 注意事项

- 请确保 `slides` 文件夹与 `presentation.html` 在同一目录下
- 建议使用现代浏览器（Chrome、Firefox、Edge）打开
- 如需编辑单个幻灯片，可直接修改 `slides` 文件夹中的 HTML 文件
"""
    
    def export(
        self,
        output_path: Optional[Path] = None,
        include_html: bool = True,
        include_pptx: bool = True,
        include_screenshots: bool = True,
        include_plan: bool = True,
        include_speech: bool = True,
        return_bytes: bool = False
    ) -> Union[Path, bytes]:
        """Export all assets as a ZIP file.
        
        Args:
            output_path: Output ZIP file path
            include_html: Include HTML slides and combined presentation
            include_pptx: Include PPTX file (requires screenshots)
            include_screenshots: Include screenshots
            include_plan: Include presentation plan
            include_speech: Include speech script and coaching
            return_bytes: If True, return ZIP as bytes instead of saving to file
            
        Returns:
            Path to exported ZIP file, or bytes if return_bytes=True
        """
        from export.html_exporter import HTMLExporter
        import tempfile
        
        # Read manifest for title
        manifest = self._get_manifest()
        title = "Presentation"
        if manifest:
            title = manifest.get("title", "Presentation")
        
        # Create a temporary directory for building the package
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            
            # Generate presentation.html with correct paths for ZIP structure
            html_path = None
            if include_html:
                try:
                    html_exporter = HTMLExporter(self.workspace_path)
                    html_path = tmp_path / "presentation.html"
                    html_exporter.export_combined(
                        output_path=html_path,
                        title=title,
                        slides_path_prefix="slides/"
                    )
                except Exception as e:
                    logger.warning(f"Failed to generate HTML: {e}")
            
            # Generate PPTX if requested
            pptx_path = None
            screenshot_paths = []
            if include_pptx or include_screenshots:
                try:
                    logger.info("Generating screenshots for PPTX...")
                    screenshot_paths = self._screenshot_slides()
                    
                    if screenshot_paths and include_pptx:
                        logger.info("Creating PPTX from screenshots...")
                        pptx_path = tmp_path / "presentation.pptx"
                        self._create_pptx_from_screenshots(
                            screenshot_paths,
                            output_path=pptx_path,
                            title=title
                        )
                except Exception as e:
                    logger.warning(f"Failed to generate PPTX: {e}")
                    import traceback
                    traceback.print_exc()
            
            # Generate speech materials if requested
            speech_script = None
            speech_coaching = None
            if include_speech:
                try:
                    logger.info("Generating speech materials...")
                    speech_script, speech_coaching = self._generate_speech_materials()
                except Exception as e:
                    logger.warning(f"Failed to generate speech materials: {e}")
            
            # Generate README content
            readme_content = self._generate_readme(
                title=title,
                has_pptx=pptx_path is not None and pptx_path.exists() if pptx_path else False,
                has_speech=speech_script is not None,
                has_coaching=speech_coaching is not None
            )
            
            # Create ZIP
            if return_bytes:
                # Create ZIP in memory
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                    self._write_zip_contents(
                        zf,
                        html_path=html_path,
                        pptx_path=pptx_path,
                        speech_script=speech_script,
                        speech_coaching=speech_coaching,
                        readme_content=readme_content,
                        include_screenshots=include_screenshots,
                        screenshot_paths=screenshot_paths,
                        include_plan=include_plan
                    )
                
                zip_buffer.seek(0)
                logger.info("Created ZIP package in memory")
                return zip_buffer.getvalue()
            else:
                # Create ZIP to file
                if output_path is None:
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_path = self.workspace_path / f"presentation_{timestamp}.zip"
                
                output_path = Path(output_path)
                output_path.parent.mkdir(parents=True, exist_ok=True)
                
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                    self._write_zip_contents(
                        zf,
                        html_path=html_path,
                        pptx_path=pptx_path,
                        speech_script=speech_script,
                        speech_coaching=speech_coaching,
                        readme_content=readme_content,
                        include_screenshots=include_screenshots,
                        screenshot_paths=screenshot_paths,
                        include_plan=include_plan
                    )
                
                logger.info(f"Created ZIP package: {output_path}")
                return output_path
    
    def _write_zip_contents(
        self,
        zf: zipfile.ZipFile,
        html_path: Optional[Path],
        pptx_path: Optional[Path],
        speech_script: Optional[str],
        speech_coaching: Optional[str],
        readme_content: str,
        include_screenshots: bool,
        screenshot_paths: List[Path],
        include_plan: bool
    ) -> None:
        """Write all contents to the ZIP file."""
        # Add presentation.html
        if html_path and html_path.exists():
            zf.write(html_path, 'presentation.html')
        
        # Add presentation.pptx if generated
        if pptx_path and pptx_path.exists():
            zf.write(pptx_path, 'presentation.pptx')
        
        # Add speech materials if generated
        if speech_script:
            zf.writestr('speech_script.md', speech_script)
        if speech_coaching:
            zf.writestr('speech_coaching.md', speech_coaching)
        
        # Add README.md
        zf.writestr('README.md', readme_content)
        
        # Add all files from slides directory
        if self.slides_path.exists():
            for file_path in self.slides_path.iterdir():
                if file_path.is_file():
                    arcname = f'slides/{file_path.name}'
                    zf.write(file_path, arcname)
        
        # Add screenshots if requested
        if include_screenshots and screenshot_paths:
            for screenshot_path in screenshot_paths:
                if screenshot_path.exists():
                    arcname = f'screenshots/{screenshot_path.name}'
                    zf.write(screenshot_path, arcname)
        
        # Add presentation plan if requested
        if include_plan:
            plan_path = self.slides_path / "presentation_plan.json"
            if plan_path.exists():
                zf.write(plan_path, "presentation_plan.json")
    
    def get_export_size_estimate(self) -> int:
        """Estimate total export size in bytes."""
        total = 0
        
        paths_to_check = [
            ("slides", "*.html"),
            ("slides", "*.json"),
            ("screenshots", "*.png"),
            (".", "*.pptx"),
            (".", "*.html"),
        ]
        
        for subdir, pattern in paths_to_check:
            path = self.workspace_path / subdir
            if path.exists():
                for f in path.glob(pattern):
                    total += f.stat().st_size
        
        return total